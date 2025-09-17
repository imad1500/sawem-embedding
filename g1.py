# -*- coding: utf-8 -*-
"""
Script de génération d'une présentation PowerPoint pour la formation
"Python pour l'analyse de données" - Session de 2h (9h00-11h00)

Basé sur le programme détaillé de formation et le fichier "Python 2020.pptx"
Adapté pour une présentation interactive avec notes détaillées pour le présentateur

Auteur: Formation Python 2025
Date: Octobre 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

# Créer une nouvelle présentation
prs = Presentation()

# --- Fonctions utilitaires pour la création des diapositives ---
def add_title_slide(prs, title, subtitle="", notes=""):
    """Ajoute une diapositive de titre avec mise en forme."""
    slide_layout = prs.slide_layouts[0]  # Layout titre
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    if subtitle:
        slide.placeholders[1].text = subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def add_section_slide(prs, title, subtitle="", notes=""):
    """Ajoute une diapositive de section pour marquer une nouvelle partie."""
    slide_layout = prs.slide_layouts[2]  # Layout section
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def add_content_slide(prs, title, content, notes=""):
    """Ajoute une diapositive avec un titre et du contenu textuel."""
    slide_layout = prs.slide_layouts[1]  # Layout contenu
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    if content:
        slide.placeholders[1].text = content
        content_shape = slide.placeholders[1]
        content_shape.text_frame.paragraphs[0].font.size = Pt(20)
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def add_code_slide(prs, title, code, notes="", explanation=""):
    """Ajoute une diapositive pour un exemple de code."""
    slide_layout = prs.slide_layouts[6]  # Layout vierge
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    
    # Explication si fournie
    if explanation:
        exp_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1))
        exp_frame = exp_shape.text_frame
        exp_para = exp_frame.paragraphs[0]
        exp_para.text = explanation
        exp_para.font.size = Pt(16)
        code_top = Inches(2.2)
    else:
        code_top = Inches(1.5)
    
    # Zone de code avec fond coloré
    code_box = slide.shapes.add_textbox(Inches(0.5), code_top, Inches(9), Inches(4.5))
    text_frame = code_box.text_frame
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_top = Inches(0.1)
    
    p = text_frame.paragraphs[0]
    p.text = code
    p.font.name = 'Courier New'
    p.font.size = Pt(14)
    
    # Ajouter fond coloré au code
    fill = code_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def add_two_column_slide(prs, title, left_content, right_content, notes=""):
    """Ajoute une diapositive avec deux colonnes."""
    slide_layout = prs.slide_layouts[6]  # Layout vierge
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    
    # Colonne gauche
    left_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_shape.text_frame
    left_frame.text = left_content
    for para in left_frame.paragraphs:
        para.font.size = Pt(16)
    
    # Colonne droite
    right_shape = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_shape.text_frame
    right_frame.text = right_content
    for para in right_frame.paragraphs:
        para.font.size = Pt(16)
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

# =============================================================================
# DÉBUT DE LA PRÉSENTATION - 6 DIAPOSITIVES D'INTRODUCTION À LA FORMATION
# =============================================================================

print("Génération de la présentation Python - Formation 2h...")

# Diapositive 1 : Page de titre principale
add_title_slide(
    prs,
    "PYTHON POUR L'ANALYSE DE DONNÉES",
    "Formation Jour 1 - Fondamentaux\n"
    "Durée : 2h (9h00 - 11h00)\n"
    "Dates : 6-7, 13-14, 20-21 octobre 2025\n\n"
    "Formateur : [Votre nom]\n"
    "Contact : [votre.email@domaine.fr]",
    notes="""NOTES PRÉSENTATEUR :
    
    ACCUEIL ET MISE EN CONTEXTE (10 min) :
    - Souhaiter la bienvenue à tous les participants
    - Se présenter brièvement (nom, expérience, domaines d'expertise)
    - Expliquer le contexte de cette formation sur 4 jours
    - Mentionner que cette session de 2h pose les fondements
    
    POINTS CLÉS À MENTIONNER :
    - Cette formation s'adresse aux développeurs, data analysts, scientifiques
    - Prérequis : connaissances de base en Python
    - Approche très pratique avec 70% de pratique, 30% théorie
    - Outils : Jupyter Notebook, PyCharm/VS Code, bibliothèques principales
    
    CRÉER L'AMBIANCE :
    - Formation interactive, encourager les questions
    - L'objectif est que chacun reparte avec des bases solides
    - Nous couvrirons beaucoup de terrain en 2h, donc rythme soutenu"""
)

# Diapositive 2 : Tour de table et présentation des participants
add_content_slide(
    prs,
    "Faisons connaissance",
    "TOUR DE TABLE (10 minutes)\n\n"
    "Pour chaque participant :\n"
    "• Nom et fonction\n"
    "• Entreprise/organisation\n"
    "• Expérience actuelle avec Python\n"
    "• Objectifs et attentes pour cette formation\n"
    "• Domaines d'application visés\n\n"
    "💡 Cela nous permettra d'adapter les exemples à vos besoins !",
    notes="""NOTES PRÉSENTATEUR :
    
    GESTION DU TOUR DE TABLE :
    - Limiter chaque présentation à 1-2 minutes maximum
    - Prendre des notes sur les profils pour adapter les exemples
    - Identifier les niveaux : débutants, intermédiaires, avancés
    - Noter les domaines d'application mentionnés
    
    QUESTIONS À POSER SI NÉCESSAIRE :
    - "Avez-vous déjà utilisé Pandas/NumPy ?"
    - "Travaillez-vous avec des données ? Quel type ?"
    - "Quels sont vos outils actuels d'analyse ?"
    
    ADAPTER LE CONTENU :
    - Si beaucoup de débutants → insister sur les bases
    - Si profils avancés → aller plus vite sur la syntaxe
    - Noter les cas d'usage pour les exemples"""
)

# Diapositive 3 : Vue d'ensemble de la formation complète
add_content_slide(
    prs,
    "Programme de formation - Vue d'ensemble",
    "📅 JOUR 1 (aujourd'hui) : Fondamentaux de Python\n"
    "   • Historique, syntaxe, types de données\n"
    "   • Structures de contrôle, fonctions, POO\n"
    "   • Gestion des fichiers (après-midi)\n\n"
    "📅 JOUR 2 : Manipulation et Visualisation\n"
    "   • NumPy, Pandas, Matplotlib\n"
    "   • Gestion des exceptions\n\n"
    "📅 JOUR 3 : Bases de données et Machine Learning\n"
    "   • SQL, serveurs web, TensorFlow/Keras\n\n"
    "📅 JOUR 4 : Optimisation et évaluations\n"
    "   • Programmation asynchrone, projets pratiques",
    notes="""NOTES PRÉSENTATEUR :
    
    VISION GLOBALE :
    - Expliquer la progression pédagogique sur 4 jours
    - Insister sur l'approche en spirale : on revient sur les concepts
    - Chaque jour s'appuie sur le précédent
    
    AUJOURD'HUI SPÉCIFIQUEMENT :
    - 2h de théorie intensive sur les fondamentaux
    - L'après-midi sera consacré aux ateliers pratiques
    - Important de bien maîtriser les bases pour la suite
    
    RASSURER LES PARTICIPANTS :
    - Rythme intensif mais progression logique
    - Beaucoup de pratique dès cet après-midi
    - Support et ressources disponibles"""
)

# Diapositive 4 : Objectifs pédagogiques de la matinée
add_content_slide(
    prs,
    "Objectifs de cette session (9h00-11h00)",
    "🎯 À la fin de ces 2 heures, vous saurez :\n\n"
    "✅ Pourquoi Python est devenu incontournable en data science\n"
    "✅ Maîtriser la syntaxe fondamentale de Python\n"
    "✅ Manipuler les types de données essentiels\n"
    "✅ Utiliser les structures de contrôle (if, for, while)\n"
    "✅ Créer et utiliser des fonctions\n"
    "✅ Comprendre les bases de la POO\n"
    "✅ Organiser son code avec les modules\n\n"
    "📋 Base solide pour les ateliers de cet après-midi !",
    notes="""NOTES PRÉSENTATEUR :
    
    OBJECTIFS PÉDAGOGIQUES CLAIRS :
    - Expliquer que ces objectifs sont mesurables
    - À la fin, faire un récapitulatif pour vérifier l'atteinte
    - Ces bases sont essentielles pour NumPy/Pandas
    
    GESTION DES ATTENTES :
    - 2h c'est court pour tout voir en détail
    - On privilégie la compréhension globale
    - L'approfondissement se fera en pratique
    
    LIEN AVEC LA SUITE :
    - Expliquer pourquoi chaque point est important
    - Les fonctions → pour structurer les analyses
    - La POO → pour comprendre Pandas
    - Les modules → pour les bibliothèques scientifiques"""
)

# Diapositive 5 : Méthodologie et outils
add_two_column_slide(
    prs,
    "Méthodologie et outils de travail",
    "📚 MÉTHODOLOGIE :\n\n"
    "• 70% pratique, 30% théorie\n"
    "• Ateliers guidés pas à pas\n"
    "• Exemples concrets et applicables\n"
    "• Questions encouragées à tout moment\n"
    "• Apprentissage par l'erreur\n\n"
    "⚡ RYTHME SOUTENU :\n"
    "• Beaucoup de contenu en 2h\n"
    "• Focus sur l'essentiel\n"
    "• Approfondissement en TP",
    "🛠️ OUTILS UTILISÉS :\n\n"
    "• Python 3.8+ (Anaconda/Miniconda)\n"
    "• Jupyter Notebook\n"
    "• PyCharm ou VS Code\n"
    "• Bibliothèques : NumPy, Pandas, Matplotlib\n\n"
    "💻 ENVIRONNEMENT :\n"
    "• Installation cet après-midi (11h-12h30)\n"
    "• Environnements virtuels\n"
    "• Gestion des dépendances avec pip/conda",
    notes="""NOTES PRÉSENTATEUR :
    
    RASSURER SUR LA MÉTHODOLOGIE :
    - Pas de cours magistral passif
    - Interaction permanente souhaitée
    - Droit à l'erreur, c'est formateur
    
    EXPLIQUER LE RYTHME :
    - 2h de contenu dense mais structuré
    - Pauses courtes si nécessaire
    - L'installation sera faite ensemble cet après-midi
    
    OUTILS :
    - Ne pas s'inquiéter si pas encore installé
    - On utilisera l'interpréteur Python de base ce matin
    - Jupyter/IDE cet après-midi pour les TP"""
)

# Diapositive 6 : Plan détaillé de la matinée
add_content_slide(
    prs,
    "Plan détaillé de la matinée",
    "⏰ 9h00-9h15 : Accueil et présentation (15min)\n\n"
    "⏰ 9h15-9h45 : Historique et intérêt de Python (30min)\n"
    "   • Pourquoi Python ? Écosystème, communauté\n"
    "   • Concurrents et positionnement\n\n"
    "⏰ 9h45-10h30 : Syntaxe et types de données (45min)\n"
    "   • Variables, types simples et complexes\n"
    "   • Listes, tuples, dictionnaires\n\n"
    "⏰ 10h30-11h00 : Structures et concepts avancés (30min)\n"
    "   • Conditions, boucles, fonctions\n"
    "   • Introduction POO et modules\n\n"
    "🎯 Questions/réponses intégrées tout au long !",
    notes="""NOTES PRÉSENTATEUR :
    
    GESTION DU TEMPS :
    - Respecter scrupuleusement les créneaux
    - Prévoir une marge sur chaque section
    - Si retard, prioriser les concepts essentiels
    
    TRANSITIONS :
    - Annoncer clairement les transitions
    - Faire des liens entre les parties
    - Rappeler où on en est dans le programme
    
    INTERACTION :
    - Poser des questions pour vérifier la compréhension
    - Inviter aux questions à tout moment
    - Utiliser des exemples du tour de table"""
)

# =============================================================================
# COURS PRINCIPAL - 45 DIAPOSITIVES SUR PYTHON
# =============================================================================

# SECTION 1 : HISTORIQUE ET INTÉRÊT (9h15-9h45)
add_section_slide(
    prs, 
    "HISTORIQUE ET INTÉRÊT DE PYTHON",
    "Pourquoi Python est-il devenu le langage de référence en data science ?",
    notes="""SECTION 1 - 30 MINUTES
    
    OBJECTIFS DE CETTE SECTION :
    - Comprendre les raisons du succès de Python
    - Situer Python par rapport à la concurrence
    - Découvrir l'écosystème et la communauté
    
    APPROCHE PÉDAGOGIQUE :
    - Partir de l'expérience des participants
    - Montrer des chiffres concrets
    - Exemples d'entreprises utilisatrices"""
)

# Diapositive 7 : Naissance et évolution de Python
add_content_slide(
    prs,
    "Naissance et évolution de Python",
    "👨‍💻 CRÉATEUR : Guido van Rossum (1989-1991)\n"
    "   • Développé aux Pays-Bas (CWI)\n"
    "   • Nom inspiré des 'Monty Python'\n"
    "   • 'Benevolent Dictator For Life' jusqu'en 2018\n\n"
    "📈 ÉVOLUTION DES VERSIONS :\n"
    "   • Python 1.0 (1994) - Premières bases\n"
    "   • Python 2.0 (2000) - Listes en compréhension\n"
    "   • Python 3.0 (2008) - Rupture de compatibilité\n"
    "   • Python 3.9+ (2020+) - Performances optimisées\n\n"
    "🏆 Aujourd'hui : Langage n°1 selon TIOBE Index !",
    notes="""NOTES PRÉSENTATEUR :
    
    ANECDOTES INTÉRESSANTES :
    - Guido van Rossum travaille maintenant chez Microsoft
    - Le nom vient du groupe comique, pas du serpent
    - Python 2 vs 3 : grande migration achevée en 2020
    
    INTERAGIR AVEC L'AUDIENCE :
    - "Qui a déjà utilisé Python 2 ?"
    - "Quelles versions utilisez-vous actuellement ?"
    
    CHIFFRES MARQUANTS :
    - +200% de croissance d'usage depuis 2017
    - Language of the year 2021 selon TIOBE"""
)

# Diapositive 8 : Caractéristiques techniques clés
add_two_column_slide(
    prs,
    "Caractéristiques techniques de Python",
    "✅ AVANTAGES :\n\n"
    "• Langage INTERPRÉTÉ\n"
    "  → Pas de compilation\n"
    "  → Test immédiat du code\n\n"
    "• PORTABILITÉ totale\n"
    "  → Windows, Mac, Linux\n"
    "  → Même code partout\n\n"
    "• SYNTAXE SIMPLE\n"
    "  → Lisible comme l'anglais\n"
    "  → Moins d'erreurs\n\n"
    "• TYPAGE DYNAMIQUE\n"
    "  → Flexibilité maximale\n"
    "  → Productivité accrue",
    "⚠️ INCONVÉNIENTS :\n\n"
    "• VITESSE D'EXÉCUTION\n"
    "  → Plus lent que C/C++\n"
    "  → Compensé par les bibliothèques\n\n"
    "• GIL (Global Interpreter Lock)\n"
    "  → Limite le multithreading\n"
    "  → Solutions : multiprocessing\n\n"
    "• CONSOMMATION MÉMOIRE\n"
    "  → Plus gourmand\n"
    "  → Non critique sur machines modernes\n\n"
    "🎯 Bilan : Avantages >> Inconvénients\n"
    "pour l'analyse de données !",
    notes="""NOTES PRÉSENTATEUR :
    
    EXPLIQUER LES CONCEPTS TECHNIQUES :
    - Interprété vs compilé : montrer avec un exemple simple
    - GIL : concept avancé, ne pas trop détailler maintenant
    - Vitesse : NumPy/Pandas utilisent du C optimisé
    
    RASSURER SUR LES INCONVÉNIENTS :
    - La vitesse n'est pas critique pour l'analyse de données
    - Les bibliothèques scientifiques sont très optimisées
    - La productivité compense largement
    
    EXEMPLES CONCRETS :
    - Instagram, YouTube utilisent Python massivement
    - NASA, CERN l'utilisent pour les calculs scientifiques"""
)

# Diapositive 9 : Python pour la data science - Pourquoi ce choix ?
add_content_slide(
    prs,
    "Pourquoi Python pour la Data Science ?",
    "🚀 PRODUCTIVITÉ EXCEPTIONNELLE :\n"
    "   • Code 5x plus concis que Java\n"
    "   • Prototypage rapide\n"
    "   • Debug facilité\n\n"
    "📚 ÉCOSYSTÈME RICHE :\n"
    "   • NumPy : calcul scientifique\n"
    "   • Pandas : manipulation de données\n"
    "   • Matplotlib/Seaborn : visualisation\n"
    "   • Scikit-learn : machine learning\n"
    "   • TensorFlow/PyTorch : deep learning\n\n"
    "👥 COMMUNAUTÉ ACTIVE :\n"
    "   • 400,000+ packages sur PyPI\n"
    "   • Documentation excellente\n"
    "   • Support communautaire",
    notes="""NOTES PRÉSENTATEUR :
    
    EXPLIQUER LA PRODUCTIVITÉ :
    - Montrer un exemple simple de code Python vs autre langage
    - Expliquer la philosophie "batteries included"
    - Temps de développement divisé par 3-4 en moyenne
    
    ÉCOSYSTÈME :
    - Ces bibliothèques seront vues dans les jours suivants
    - Chacune est un standard dans son domaine
    - Interopérabilité excellente entre elles
    
    COMMUNAUTÉ :
    - Stack Overflow : Python dans le top 3 des questions
    - PyPI : plus grand dépôt de packages au monde
    - Conférences : PyCon dans chaque pays"""
)

# Diapositive 10 : La concurrence - Comparaison objective
add_content_slide(
    prs,
    "Python face à la concurrence",
    "🔬 MATLAB :\n"
    "   ✅ Très performant pour le calcul numérique\n"
    "   ❌ Propriétaire, coûteux, syntaxe spécifique\n\n"
    "📊 R :\n"
    "   ✅ Excellence en statistiques\n"
    "   ❌ Courbe d'apprentissage, moins polyvalent\n\n"
    "☕ JAVA :\n"
    "   ✅ Performance, robustesse enterprise\n"
    "   ❌ Verbosité, complexité, moins adapté à la data\n\n"
    "⚡ JULIA :\n"
    "   ✅ Performance native, syntaxe moderne\n"
    "   ❌ Écosystème encore limité, communauté réduite",
    notes="""NOTES PRÉSENTATEUR :
    
    RESTER OBJECTIF :
    - Chaque langage a ses domaines d'excellence
    - Python n'est pas parfait partout
    - Le choix dépend du contexte et des équipes
    
    EXPÉRIENCE DES PARTICIPANTS :
    - "Qui a déjà utilisé R/MATLAB/Java ?"
    - Recueillir les retours d'expérience
    - Expliquer les migrations vers Python
    
    TENDANCES ACTUELLES :
    - Python gagne des parts de marché chaque année
    - Migration progressive depuis R et MATLAB
    - Choix par défaut des nouvelles équipes data"""
)

# Diapositive 11 : L'écosystème Python - Vue d'ensemble
add_two_column_slide(
    prs,
    "L'écosystème Python pour la Data Science",
    "🧮 CALCUL SCIENTIFIQUE :\n"
    "• NumPy - Arrays multidimensionnels\n"
    "• SciPy - Algorithmes scientifiques\n"
    "• SymPy - Calcul symbolique\n\n"
    "📊 MANIPULATION DE DONNÉES :\n"
    "• Pandas - DataFrames et séries\n"
    "• Polars - Alternative haute performance\n"
    "• Dask - Calcul parallèle\n\n"
    "📈 VISUALISATION :\n"
    "• Matplotlib - Graphiques de base\n"
    "• Seaborn - Visualisation statistique\n"
    "• Plotly - Graphiques interactifs",
    "🤖 MACHINE LEARNING :\n"
    "• Scikit-learn - ML classique\n"
    "• TensorFlow - Deep Learning Google\n"
    "• PyTorch - Deep Learning Facebook\n"
    "• Keras - Interface haut niveau\n\n"
    "🔧 OUTILS COMPLÉMENTAIRES :\n"
    "• Jupyter - Notebooks interactifs\n"
    "• Anaconda - Distribution scientifique\n"
    "• Streamlit - Applications web\n"
    "• FastAPI - APIs modernes\n\n"
    "💾 DONNÉES :\n"
    "• SQLAlchemy - Base de données\n"
    "• Requests - APIs REST",
    notes="""NOTES PRÉSENTATEUR :
    
    NE PAS TOUT DÉTAILLER MAINTENANT :
    - Vue d'ensemble pour montrer la richesse
    - Nous verrons NumPy/Pandas/Matplotlib en détail demain
    - Les autres selon les besoins des participants
    
    EXPLIQUER LA COMPLÉMENTARITÉ :
    - Ces outils se combinent naturellement
    - Pandas utilise NumPy, Matplotlib utilise NumPy, etc.
    - Écosystème cohérent et intégré
    
    ÉVOLUTION RAPIDE :
    - Nouveaux outils régulièrement
    - Communauté très active
    - Importance de rester à jour"""
)

# Diapositive 12 : Classement et adoption industrielle
add_content_slide(
    prs,
    "Python dans l'industrie - Chiffres clés 2024",
    "📈 CLASSEMENTS :\n"
    "   • TIOBE Index : #1 mondial\n"
    "   • IEEE Spectrum : #1 pour 4ème année\n"
    "   • Stack Overflow : #3 le plus aimé\n"
    "   • GitHub : #2 en nombre de projets\n\n"
    "🏢 ADOPTIONS NOTABLES :\n"
    "   • Netflix → Systèmes de recommandation\n"
    "   • Instagram → Backend principal\n"
    "   • NASA → Analyse de données spatiales\n"
    "   • Spotify → Analytics et ML\n"
    "   • Uber → Tarification dynamique\n\n"
    "💼 MARCHÉ DE L'EMPLOI :\n"
    "   • +35% d'offres Python en 2 ans\n"
    "   • Salaires moyens : 55-75k€ en France",
    notes="""NOTES PRÉSENTATEUR :
    
    CHIFFRES RÉCENTS :
    - Mettre à jour avec les dernières données disponibles
    - Montrer la progression constante
    - Expliquer les méthodologies de classement
    
    CAS D'USAGE INSPIRANTS :
    - Netflix : algorithmes de recommandation pour 200M+ users
    - Instagram : gère des milliards de photos
    - NASA : traitement d'images satellite et télescope
    
    OPPORTUNITÉS PROFESSIONNELLES :
    - Marché en forte croissance
    - Compétence très recherchée
    - Polyvalence : dev, data, ML, ops"""
)

# SECTION 2 : SYNTAXE ET TYPES DE DONNÉES (9h45-10h30)
add_section_slide(
    prs, 
    "SYNTAXE DE BASE ET TYPES DE DONNÉES",
    "Les fondements du langage Python",
    notes="""SECTION 2 - 45 MINUTES
    
    PARTIE LA PLUS IMPORTANTE :
    - 50% du temps sur cette section
    - Concepts essentiels pour la suite
    - Beaucoup d'exemples pratiques
    
    GESTION DU TEMPS :
    - Variables et types simples : 15 min
    - Structures de données : 20 min  
    - Opérateurs et expressions : 10 min
    
    INTERACTIVITÉ :
    - Montrer les exemples en direct si possible
    - Faire participer avec des questions simples
    - Encourager à noter les concepts nouveaux"""
)

# Diapositive 13 : Premiers pas - Variables et affectation
add_code_slide(
    prs,
    "Premiers pas : Variables et affectation",
    """# En Python, pas besoin de déclarer le type !
nom = "Alice"           # String (chaîne de caractères)
age = 30               # Integer (entier)
taille = 1.65          # Float (nombre à virgule)
est_majeur = True      # Boolean (booléen)

# Vérification des types
print(type(nom))       # <class 'str'>
print(type(age))       # <class 'int'>
print(type(taille))    # <class 'float'>
print(type(est_majeur)) # <class 'bool'>""",
    explanation="Le typage dynamique : Python détermine automatiquement le type",
    notes="""NOTES PRÉSENTATEUR :

CONCEPTS CLÉS À EXPLIQUER :
- Typage dynamique vs statique (Java, C++)
- Pas de déclaration de type nécessaire
- Type déterminé à l'exécution
- Fonction type() pour vérifier

DÉMONSTRATION INTERACTIVE :
- Taper ces exemples en direct dans un interpréteur
- Montrer le résultat de print(type(...))
- Expliquer les noms de classes Python

ATTENTION AUX DÉBUTANTS :
- Variables sensibles à la casse : Age ≠ age
- Conventions de nommage : snake_case recommandé
- Mots réservés interdits comme noms de variables"""
)

# Diapositive 14 : Affectations multiples et échange de variables
add_code_slide(
    prs,
    "Affectations multiples et échanges",
    """# Affectations multiples
a = b = c = 5          # Même valeur à toutes
print(a, b, c)         # 5 5 5

# Affectation parallèle (unpacking)
x, y, z = 10, 20, 30   # Assigne en parallèle
print(x, y, z)         # 10 20 30

# Le fameux échange de variables Python !
a = 100
b = 200
print("Avant:", a, b)   # Avant: 100 200

# Échange en une seule ligne (magique!)
a, b = b, a
print("Après:", a, b)   # Après: 200 100""",
    explanation="Python permet des affectations élégantes et concises",
    notes="""NOTES PRÉSENTATEUR :

ÉMERVEILLER L'AUDIENCE :
- L'échange de variables en une ligne est iconique Python
- Dans d'autres langages : 3 lignes avec variable temporaire
- Montrer la différence avec Java/C++

CONCEPT TECHNIQUE :
- Unpacking = déballage de séquence
- Côté droit évalué en premier (tuple temporaire)
- Très utilisé en data science pour les coordonnées

EXERCICE MENTAL :
- Faire deviner le résultat avant d'exécuter
- Expliquer pourquoi ça marche
- Applications pratiques"""
)

# Diapositive 15 : Opérateurs arithmétiques et comparaisons
add_code_slide(
    prs,
    "Opérateurs essentiels",
    """# Opérateurs arithmétiques
a, b = 10, 3
print(a + b)    # 13 - Addition
print(a - b)    # 7  - Soustraction  
print(a * b)    # 30 - Multiplication
print(a / b)    # 3.333... - Division (toujours float en Python 3)
print(a // b)   # 3  - Division entière
print(a % b)    # 1  - Modulo (reste)
print(a ** b)   # 1000 - Puissance

# Comparaisons (résultat = booléen)
print(a > b)    # True
print(a == b)   # False
print(a != b)   # True

# Opérateurs logiques
print(True and False)  # False
print(True or False)   # True
print(not True)        # False""",
    explanation="Attention à la division : / donne toujours un float, // pour l'entier",
    notes="""NOTES PRÉSENTATEUR :

PIÈGE CLASSIQUE Python 2 vs 3 :
- Python 2 : 10/3 = 3 (division entière)
- Python 3 : 10/3 = 3.333... (division flottante)
- Toujours utiliser // pour division entière

UTILITÉ DU MODULO :
- Vérifier parité : x % 2 == 0
- Opérations cycliques (horaires, calendriers)
- Algorithmes de hachage

OPÉRATEURS LOGIQUES :
- and, or, not (mots anglais, pas symboles)
- Court-circuit : and s'arrête au premier False
- Différent de & | (opérateurs bit à bit)"""
)

# Diapositive 16 : Chaînes de caractères - Manipulation de base
add_code_slide(
    prs,
    "Chaînes de caractères (Strings)",
    """# Création de chaînes
nom = 'Alice'           # Guillemets simples
message = "Bonjour!"    # Guillemets doubles  
long_text = """Texte
sur plusieurs
lignes"""               # Triple guillemets

# Indexation et slicing
prenom = "Jean-Michel"
print(prenom[0])        # 'J' - Premier caractère
print(prenom[-1])       # 'l' - Dernier caractère
print(prenom[0:4])      # 'Jean' - Slice début:fin
print(prenom[5:])       # 'Michel' - Slice du 5ème à la fin

# Opérations sur chaînes
print(len(prenom))      # 11 - Longueur
print("jean" in prenom.lower())  # True - Recherche
print(prenom.upper())   # 'JEAN-MICHEL' - Majuscules""",
    explanation="Les strings sont IMMUABLES : on ne peut pas les modifier directement",
    notes="""NOTES PRÉSENTATEUR :

CONCEPT D'IMMUTABILITÉ :
- prenom[0] = 'P' → ERREUR !
- Il faut créer une nouvelle chaîne
- Optimisation mémoire de Python

INDEXATION NÉGATIVE :
- Très utile en Python : -1 = dernier, -2 = avant-dernier
- Évite de calculer len(chaine) - 1

SLICING AVANCÉ :
- [début:fin:pas] : prenom[::2] tous les 2 caractères
- [::-1] pour inverser une chaîne
- Très utilisé en manipulation de données"""
)

# Diapositive 17 : Formatage de chaînes moderne
add_code_slide(
    prs,
    "Formatage de chaînes - f-strings (Python 3.6+)",
    """nom = "Marie"
age = 28
taille = 1.68

# Ancienne méthode (à éviter)
message1 = "Je m'appelle " + nom + " et j'ai " + str(age) + " ans"

# Méthode format() (acceptable)
message2 = "Je m'appelle {} et j'ai {} ans".format(nom, age)

# f-strings (recommandée !) ⭐
message3 = f"Je m'appelle {nom} et j'ai {age} ans"
print(message3)  # Je m'appelle Marie et j'ai 28 ans

# f-strings avec expressions
print(f"Taille: {taille:.1f}m")           # Taille: 1.7m
print(f"Dans 5 ans: {age + 5} ans")       # Dans 5 ans: 33 ans
print(f"Majeure: {age >= 18}")            # Majeure: True""",
    explanation="Les f-strings : plus lisibles, plus rapides, plus puissantes !",
    notes="""NOTES PRÉSENTATEUR :

ÉVOLUTION DU FORMATAGE :
- Concaténation (+) : lente et peu lisible
- .format() : amélioration mais verbeux
- f-strings : révolution depuis Python 3.6

AVANTAGES f-strings :
- Plus rapides à l'exécution
- Plus lisibles et maintenables
- Expressions directement intégrées
- Formatage puissant intégré

FORMATAGE AVANCÉ :
- {valeur:.2f} : 2 décimales
- {valeur:>10} : aligné à droite sur 10 caractères
- {valeur:,} : séparateurs de milliers"""
)

# Diapositive 18 : Listes - Le type de données fondamental
add_code_slide(
    prs,
    "Listes : La structure de données fondamentale",
    """# Création et manipulation de base
nombres = [1, 5, 3, 9, 2]
print(nombres)              # [1, 5, 3, 9, 2]

# Listes hétérogènes (types mélangés)
mixte = ["Alice", 25, True, 3.14]
print(mixte)               # ['Alice', 25, True, 3.14]

# Accès aux éléments (comme les strings)
print(nombres[0])          # 1 - Premier élément
print(nombres[-1])         # 2 - Dernier élément
print(nombres[1:4])        # [5, 3, 9] - Slice

# Modification (MUTABLE !)
nombres[0] = 100
print(nombres)             # [100, 5, 3, 9, 2]

# Longueur
print(len(nombres))        # 5""",
    explanation="Contrairement aux strings, les listes sont MUTABLES (modifiables)",
    notes="""NOTES PRÉSENTATEUR :

DIFFÉRENCE CRUCIALE :
- Strings : immutables
- Listes : mutables
- Conséquences importantes pour les fonctions

LISTES HÉTÉROGÈNES :
- Spécificité Python (pas toujours possible ailleurs)
- Utile mais attention aux types en data science
- NumPy préférera les types homogènes

INDEXATION ET SLICING :
- Mêmes règles que les strings
- Très utilisé en analyse de données
- Base pour comprendre NumPy arrays"""
)

# Diapositive 19 : Méthodes des listes
add_code_slide(
    prs,
    "Méthodes essentielles des listes",
    """ma_liste = [3, 1, 4, 1, 5]
print("Liste initiale:", ma_liste)

# Ajouter des éléments
ma_liste.append(9)              # Ajoute à la fin
print("Après append(9):", ma_liste)

ma_liste.insert(2, 2)           # Insert 2 à l'index 2
print("Après insert(2,2):", ma_liste)

# Supprimer des éléments  
ma_liste.remove(1)              # Supprime la première occurrence de 1
print("Après remove(1):", ma_liste)

element = ma_liste.pop()        # Supprime et retourne le dernier
print(f"Élément supprimé: {element}, Liste: {ma_liste}")

# Trier et organiser
ma_liste.sort()                 # Trie sur place
print("Après sort():", ma_liste)

ma_liste.reverse()              # Inverse l'ordre
print("Après reverse():", ma_liste)""",
    explanation="Attention : ces méthodes modifient la liste originale !",
    notes="""NOTES PRÉSENTATEUR :

MÉTHODES QUI MODIFIENT :
- append(), insert(), remove(), pop(), sort(), reverse()
- Modifient la liste originale (pas de return)
- Différent de sorted() qui retourne une nouvelle liste

DIFFÉRENCE IMPORTANTE :
- ma_liste.sort() : modifie ma_liste, retourne None
- sorted(ma_liste) : retourne nouvelle liste triée

CAS D'USAGE :
- append() : construction de liste en boucle
- remove() : nettoyage de données
- sort() : préparation pour analyse"""
)

# Diapositive 20 : Tuples - Les listes immuables
add_code_slide(
    prs,
    "Tuples : Listes immuables",
    """# Création de tuples
point = (10, 20)              # Coordonnées x, y
couleur = (255, 0, 128)       # RGB
info = ("Alice", 30, "Paris") # Nom, âge, ville

# Parenthèses optionnelles (mais recommandées)
dimensions = 1920, 1080       # Résolution écran
print(type(dimensions))       # <class 'tuple'>

# Tuple à un seul élément (attention à la virgule !)
singleton = (42,)             # Virgule obligatoire
pas_tuple = (42)              # C'est juste 42 entre parenthèses

# Accès aux éléments (comme listes)
print(point[0])               # 10
print(info[1:])               # (30, 'Paris')

# IMPOSSIBLE de modifier !
# point[0] = 15               # ❌ ERREUR !

# Unpacking très utile
x, y = point
nom, age, ville = info
print(f"Coordonnées: x={x}, y={y}")""",
    explanation="Tuples = listes non modifiables, parfaites pour données fixes",
    notes="""NOTES PRÉSENTATEUR :

QUAND UTILISER LES TUPLES :
- Coordonnées, couleurs RGB
- Données qui ne doivent pas changer
- Clés de dictionnaire (immuables)
- Return de plusieurs valeurs de fonction

PIÈGE CLASSIQUE :
- (42) n'est PAS un tuple, c'est 42
- (42,) EST un tuple à un élément
- Python a besoin de la virgule

UNPACKING PUISSANT :
- x, y = point très idiomatique
- Échange de variables : a, b = b, a
- Parcours de listes de tuples"""
)

# Diapositive 21 : Dictionnaires - Clés et valeurs
add_code_slide(
    prs,
    "Dictionnaires : Associations clé-valeur",
    """# Création de dictionnaires
personne = {
    "nom": "Dupont",
    "prenom": "Jean", 
    "age": 35,
    "ville": "Lyon"
}

# Accès aux valeurs
print(personne["nom"])        # "Dupont"
print(personne.get("age"))    # 35
print(personne.get("pays", "France"))  # "France" (valeur par défaut)

# Modification et ajout
personne["age"] = 36          # Modification
personne["email"] = "j.dupont@email.com"  # Ajout

# Parcours du dictionnaire  
for cle in personne:          # Parcours des clés
    print(f"{cle}: {personne[cle]}")
    
# Ou plus pythonique :
for cle, valeur in personne.items():
    print(f"{cle}: {valeur}")""",
    explanation="Dictionnaires = tables de hachage, accès très rapide par clé",
    notes="""NOTES PRÉSENTATEUR :

ACCÈS SÉCURISÉ :
- dict["clé"] : lève KeyError si clé inexistante
- dict.get("clé") : retourne None si inexistante
- dict.get("clé", "défaut") : retourne valeur par défaut

TYPES DE CLÉS :
- Seulement types immuables : str, int, tuple
- Pas de listes comme clés !
- Strings le plus courant

PERFORMANCE :
- Accès O(1) en moyenne
- Très efficace pour lookups
- Base de Pandas DataFrames"""
)

# Diapositive 22 : Compréhensions de listes - Python idiomatique
add_code_slide(
    prs,
    "Compréhensions de listes (List Comprehensions)",
    """# Méthode traditionnelle
carres_classique = []
for i in range(5):
    carres_classique.append(i ** 2)
print(carres_classique)       # [0, 1, 4, 9, 16]

# Compréhension de liste (pythonique !)
carres_modern = [i ** 2 for i in range(5)]
print(carres_modern)          # [0, 1, 4, 9, 16]

# Avec condition (filtrage)
pairs = [x for x in range(10) if x % 2 == 0]
print(pairs)                  # [0, 2, 4, 6, 8]

# Plus complexe
mots = ["Python", "Java", "C++", "JavaScript"]
longueurs = [len(mot) for mot in mots if len(mot) > 4]
print(longueurs)              # [6, 10] 

# Dictionnaire en compréhension
carres_dict = {x: x**2 for x in range(5)}
print(carres_dict)            # {0: 0, 1: 1, 2: 4, 3: 9, 4: 16}""",
    explanation="Style Python par excellence : concis, lisible et performant",
    notes="""NOTES PRÉSENTATEUR :

SYNTAXE GÉNÉRALE :
- [expression for item in iterable if condition]
- Plus concis et souvent plus rapide
- Très utilisé en data science

LISIBILITÉ :
- Parfois plus complexe à lire pour débutants
- Règle : si trop complexe, utiliser boucle classique
- Maximum 2-3 niveaux de compréhension

EXTENSIONS :
- Compréhensions de dictionnaires : {k: v for...}
- Compréhensions d'ensembles : {x for...}
- Expressions génératrices : (x for...)"""
)

# SECTION 3 : STRUCTURES DE CONTRÔLE (10h30-11h00)
add_section_slide(
    prs, 
    "STRUCTURES DE CONTRÔLE ET CONCEPTS AVANCÉS",
    "Conditions, boucles, fonctions et introduction à la POO",
    notes="""SECTION 3 - 30 MINUTES FINALES
    
    DERNIÈRE LIGNE DROITE :
    - Concepts essentiels rapidement
    - Focus sur la compréhension globale
    - Beaucoup d'exemples pratiques
    
    RÉPARTITION :
    - Conditions et boucles : 10 min
    - Fonctions : 10 min
    - POO et modules : 10 min"""
)

# Diapositive 23 : Structures conditionnelles
add_code_slide(
    prs,
    "Conditions : if, elif, else",
    """age = 25
situation = "étudiant"

# Structure conditionnelle complète
if age < 18:
    statut = "mineur"
    tarif = 5
elif age < 65:
    if situation == "étudiant":
        statut = "étudiant"
        tarif = 8
    else:
        statut = "actif" 
        tarif = 12
else:
    statut = "senior"
    tarif = 6

print(f"Statut: {statut}, Tarif: {tarif}€")

# Opérateur ternaire (conditionnel inline)
message = "majeur" if age >= 18 else "mineur"
print(f"Vous êtes {message}")

# Tests de vérité Python
nom = "Alice"
if nom:  # Chaîne non vide = True
    print(f"Bonjour {nom}!")""",
    explanation="L'indentation définit les blocs (pas d'accolades comme en Java/C++)",
    notes="""NOTES PRÉSENTATEUR :

INDENTATION OBLIGATOIRE :
- Pas d'accolades {}, l'indentation structure le code
- 4 espaces recommandés (PEP 8)
- Erreur IndentationError si mal indenté

CONDITIONS PYTHON :
- Pas besoin de parenthèses autour de la condition
- elif (pas else if)
- Opérateur ternaire utile pour assignations simples

TESTS DE VÉRITÉ :
- Valeurs "falsy" : False, None, 0, "", [], {}
- Tout le reste est "truthy"
- Très utilisé pour vérifications rapides"""
)

# Diapositive 24 : Boucles for et while
add_code_slide(
    prs,
    "Boucles : for et while",
    """# Boucle for - Iteration directe sur éléments
fruits = ["pomme", "banane", "orange"]
for fruit in fruits:
    print(f"J'aime les {fruit}s")

# Boucle for avec indices si nécessaire
for i, fruit in enumerate(fruits):
    print(f"{i+1}. {fruit}")

# range() pour générer des séquences
for i in range(5):              # 0 à 4
    print(f"Compteur: {i}")
    
for i in range(2, 8, 2):        # 2,4,6 (début, fin, pas)
    print(f"Pair: {i}")

# Boucle while
compteur = 0
while compteur < 3:
    print(f"Tour {compteur + 1}")
    compteur += 1

# Break et continue
for num in range(10):
    if num == 3:
        continue    # Passe à l'itération suivante
    if num == 7:
        break       # Sort de la boucle
    print(num)""",
    explanation="for pour itérer sur des collections, while pour conditions",
    notes="""NOTES PRÉSENTATEUR :

PHILOSOPHIE PYTHON :
- for pour itérer directement sur éléments
- Pas besoin d'indices dans la plupart des cas
- enumerate() quand on a besoin des indices

RANGE() TRÈS UTILISÉ :
- range(n) : 0 à n-1
- range(start, stop) : start à stop-1  
- range(start, stop, step) : avec pas

BREAK/CONTINUE :
- break : sort complètement de la boucle
- continue : passe à l'itération suivante
- Utilisé pour gestion d'erreurs, filtres"""
)

# Diapositive 25 : Fonctions - Structurer son code
add_code_slide(
    prs,
    "Fonctions : Réutiliser et structurer",
    """# Fonction simple
def saluer(nom):
    return f"Bonjour {nom} !"

message = saluer("Alice")
print(message)                  # Bonjour Alice !

# Fonction avec plusieurs paramètres et valeurs par défaut
def calculer_prix(prix_ht, tva=0.20, reduction=0):
    prix_ttc = prix_ht * (1 + tva) * (1 - reduction)
    return round(prix_ttc, 2)

# Différentes façons d'appeler
print(calculer_prix(100))                    # 120.0
print(calculer_prix(100, 0.10))             # 110.0  
print(calculer_prix(100, reduction=0.15))    # 102.0

# Fonction qui retourne plusieurs valeurs
def analyse_nombre(n):
    est_pair = n % 2 == 0
    est_positif = n > 0
    return est_pair, est_positif

pair, positif = analyse_nombre(42)
print(f"42 est pair: {pair}, positif: {positif}")""",
    explanation="Les fonctions rendent le code modulaire et réutilisable",
    notes="""NOTES PRÉSENTATEUR :

BONNES PRATIQUES :
- Noms explicites pour fonctions et paramètres
- Une fonction = une responsabilité
- Documentation avec docstrings (on verra plus tard)

PARAMÈTRES AVANCÉS :
- Valeurs par défaut très utiles
- Appel par nom de paramètre
- *args et **kwargs pour paramètres variables

RETOUR MULTIPLE :
- Vraiment un tuple qui est retourné
- Unpacking automatique à l'affectation
- Très pratique pour coordonnées, résultats multiples"""
)

# Diapositive 26 : Portée des variables (scope)
add_code_slide(
    prs,
    "Portée des variables (scope)",
    """# Variables globales vs locales
compteur_global = 0

def incrementer():
    compteur_local = 1              # Variable locale
    return compteur_local + compteur_global

print(incrementer())                # 1

def modifier_global():
    global compteur_global          # Déclaration explicite
    compteur_global += 1

modifier_global()
print(compteur_global)              # 1

# Attention aux effets de bord avec objets mutables
def ajouter_element(liste, element):
    liste.append(element)           # Modifie la liste originale !
    return liste

ma_liste = [1, 2, 3]
nouvelle_liste = ajouter_element(ma_liste, 4)
print(ma_liste)                     # [1, 2, 3, 4] - Modifiée !
print(nouvelle_liste)               # [1, 2, 3, 4] - Même objet !""",
    explanation="Attention aux effets de bord : listes/dicts modifiés dans fonctions !",
    notes="""NOTES PRÉSENTATEUR :

RÈGLE IMPORTANTE :
- Variables locales masquent les globales
- global nécessaire pour modifier variable globale
- À éviter autant que possible (effet de bord)

OBJETS MUTABLES :
- listes, dictionnaires passés par référence
- Modifications visibles à l'extérieur de la fonction
- Source de bugs fréquents pour débutants

BONNES PRATIQUES :
- Éviter les variables globales
- Fonctions pures : même input → même output
- Retourner nouvelles valeurs plutôt que modifier"""
)

# Diapositive 27 : Introduction à la POO
add_code_slide(
    prs,
    "Introduction à la Programmation Orientée Objet",
    """# Définition d'une classe simple
class Voiture:
    # Attribut de classe (partagé)
    nb_voitures = 0
    
    # Constructeur
    def __init__(self, marque, modele, annee):
        # Attributs d'instance (propres à chaque objet)
        self.marque = marque
        self.modele = modele  
        self.annee = annee
        self.kilometrage = 0
        Voiture.nb_voitures += 1
    
    # Méthodes (comportements)
    def rouler(self, km):
        self.kilometrage += km
        return f"La {self.marque} a roulé {km} km"
    
    def __str__(self):
        return f"{self.marque} {self.modele} ({self.annee})"

# Utilisation
ma_voiture = Voiture("Toyota", "Corolla", 2020)
print(ma_voiture)                    # Toyota Corolla (2020)
print(ma_voiture.rouler(150))        # La Toyota a roulé 150 km
print(f"Total voitures: {Voiture.nb_voitures}")""",
    explanation="Classes = modèles d'objets, Objets = instances concrètes",
    notes="""NOTES PRÉSENTATEUR :

CONCEPTS FONDAMENTAUX :
- Classe = plan/modèle, Objet = instance concrète
- __init__ = constructeur (appelé à la création)
- self = référence à l'instance courante

ATTRIBUTS ET MÉTHODES :
- Attributs de classe vs attributs d'instance
- Méthodes = fonctions dans une classe
- __str__ = représentation en chaîne

POURQUOI LA POO EN DATA SCIENCE :
- Pandas DataFrame est une classe !
- Structurer code complexe
- Réutilisabilité et maintenance"""
)

# Diapositive 28 : Modules et imports
add_code_slide(
    prs,
    "Modules : Organiser et réutiliser le code",
    """# Imports de la bibliothèque standard
import math
import random
from datetime import datetime, timedelta

# Utilisation
print(math.pi)                      # 3.141592653589793
print(math.sqrt(16))               # 4.0
print(random.randint(1, 10))       # Nombre aléatoire 1-10

# Calculs de dates
aujourdhui = datetime.now()
demain = aujourdhui + timedelta(days=1)
print(f"Aujourd'hui: {aujourdhui.strftime('%d/%m/%Y')}")

# Import avec alias (très courant en data science)
import numpy as np                  # Convention universelle !
import pandas as pd                 # Convention universelle !
import matplotlib.pyplot as plt     # Convention universelle !

# Imports conditionnels et gestion d'erreurs
try:
    import tensorflow as tf
    print("TensorFlow disponible")
except ImportError:
    print("TensorFlow non installé")""",
    explanation="Modules = fichiers Python réutilisables, base de l'écosystème",
    notes="""NOTES PRÉSENTATEUR :

BIBLIOTHÈQUE STANDARD :
- Très riche : math, random, datetime, os, sys...
- Pas besoin d'installation
- Documentation excellente

CONVENTIONS DATA SCIENCE :
- np, pd, plt : conventions universelles
- Tout le monde les connaît
- Facilite lecture du code

GESTION DES DÉPENDANCES :
- try/except pour imports optionnels
- pip install pour installer packages
- Requirements.txt pour lister dépendances"""
)

# Diapositive 29 : Gestion d'erreurs - try/except (version complète)
add_code_slide(
    prs,
    "Gestion des erreurs avec try/except - Approche complète",
    """# 1. Gestion basique d'une exception spécifique
def diviser(a: float, b: float) -> float:
    try:
        resultat = a / b
        return resultat
    except ZeroDivisionError as e:
        print(f"⚠️ Erreur capturée: {type(e).__name__} - {e}")
        return float('nan')  # Retourne NaN (Not a Number) au lieu de None
    except TypeError as e:
        print(f"⚠️ Types incompatibles: {e}")
        return float('nan')

print("Division normale:", diviser(10, 2))  # 5.0
print("Division par zéro:", diviser(10, 0))  # Erreur + NaN
print("Types incompatibles:", diviser("10", "2"))  # Erreur + NaN

# 2. Structure complète try/except/else/finally
def traiter_fichier(nom_fichier: str) -> str:
    try:
        with open(nom_fichier, 'r', encoding='utf-8') as f:
            contenu = f.read()
    except FileNotFoundError:
        print(f"❌ Fichier '{nom_fichier}' introuvable")
        return ""
    except UnicodeDecodeError:
        print(f"❌ Problème d'encodage dans '{nom_fichier}'")
        return ""
    except PermissionError:
        print(f"❌ Permission refusée pour '{nom_fichier}'")
        return ""
    else:
        print(f"✅ Fichier '{nom_fichier}' lu avec succès")
        return contenu[:50] + "..." if len(contenu) > 50 else contenu
    finally:
        print(f"🔄 Opération de lecture terminée (qu'il y ait eu erreur ou non)")

# Exemple d'utilisation
resultat = traiter_fichier("exemple.txt")
print(f"Contenu extrait: {resultat}")

# 3. Création d'exceptions personnalisées
class TemperatureInvalideError(Exception):
    """Exception levée pour des températures physiquement impossibles"""
    def __init__(self, temperature, message="Température invalide"):
        self.temperature = temperature
        self.message = message
        super().__init__(f"{message}: {temperature}°C (zéro absolu = -273.15°C)")

def convertir_celsius_fahrenheit(temp_c: float) -> float:
    if temp_c < -273.15:
        raise TemperatureInvalideError(temp_c, "Température sous le zéro absolu")
    return temp_c * 9/5 + 32

# Test de l'exception personnalisée
try:
    print(convertir_celsius_fahrenheit(-300))  # Doit lever une exception
except TemperatureInvalideError as e:
    print(f"❄️ {e}")
    print(f"Température problématique: {e.temperature}°C")

# 4. Bonnes pratiques de gestion d'erreurs"""
def lire_config(config_path: str) -> dict:
    """Lit un fichier de configuration JSON avec gestion d'erreurs robuste
    import json
    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except FileNotFoundError:
        print(f"⚠️ Fichier de config '{config_path}' non trouvé. Utilisation des valeurs par défaut.")
        return {"timeout": 30, "retries": 3}
    except json.JSONDecodeError as e:
        print(f"⚠️ Fichier de config '{config_path}' mal formé: {e}. Utilisation des valeurs par défaut.")
        return {"timeout": 30, "retries": 3}
    except Exception as e:
        print(f"⚠️ Erreur inattendue lors de la lecture de '{config_path}': {e}. Utilisation des valeurs par défaut.")
        return {"timeout": 30, "retries": 3}""",
    explanation="""Gestion robuste des erreurs avec :
# 1. Capture d'exceptions spécifiques
2. Utilisation de else/finally
3. Création d'exceptions personnalisées
4. Bonnes pratiques de fallback""",
    notes="""NOTES PRÉSENTATEUR :

    STRUCTURE COMPLÈTE TRY/EXCEPT :
    - try : code à surveiller
    - except : gestion des erreurs spécifiques
    - else : exécuté si pas d'erreur
    - finally : toujours exécuté

    BONNES PRATIQUES :
    1. Capturer des exceptions SPÉCIFIQUES (pas juste Exception)
    2. Toujours fournir un message d'erreur utile
    3. Utiliser finally pour le nettoyage (fermeture de fichiers, etc.)
    4. Prévoir des valeurs de fallback quand possible
    5. Logger les erreurs pour le débogage

    EXCEPTIONS PERSONNALISÉES :
    - Héritent de Exception ou sous-classes
    - Utile pour domaine spécifique (ex: validation physique)
    - Peut inclure des données supplémentaires
    - Améliore la lisibilité du code

    CAS D'USAGE COMMUNS :
    - Validation des entrées utilisateur
    - Gestion des fichiers/IO
    - Appels réseau
    - Calculs avec contraintes physiques

    EXEMPLES CONCRETS :
    - Dans la data science : validation des jeux de données
    - Dans les APIs : gestion des erreurs HTTP
    - Dans les calculs scientifiques : vérification des domaines de validité

    À ÉVITER :
    - except: (trop large)
    - Ignorer silencieusement les erreurs (sauf cas très spécifiques)
    - Messages d'erreur vagues ("Une erreur s'est produite")"""
)

# Diapositive 30 : Concepts clés de la POO (version enrichie)
add_code_slide(
    prs,
    "Programmation Orientée Objet en Python - Concepts approfondis",
    """# 1. Définition d'une classe avec attributs et méthodes
class CompteBancaire:
    # Attribut de classe (partagé par toutes les instances)
    taux_interet = 0.01  # 1%

    def __init__(self, titulaire: str, solde: float = 0.0):
        # Attributs d'instance
        self.titulaire = titulaire
        self._solde = solde  # Convention: _ pour "protected"
        self.__historique = []  # Name mangling: __ pour "private"

    # Méthode d'instance
    def deposer(self, montant: float) -> None:
        if montant > 0:
            self._solde += montant
            self.__historique.append(f"Dépôt: +{montant}€")
        else:
            raise ValueError("Le montant doit être positif")

    def retirer(self, montant: float) -> None:
        if 0 < montant <= self._solde:
            self._solde -= montant
            self.__historique.append(f"Retrait: -{montant}€")
        else:
            raise ValueError("Montant invalide ou solde insuffisant")

    # Propriété pour accéder au solde
    @property
    def solde(self) -> float:
        return self._solde

    # Méthode de classe
    @classmethod
    def modifier_taux(cls, nouveau_taux: float) -> None:
        if 0 <= nouveau_taux <= 0.1:  # 10% max
            cls.taux_interet = nouveau_taux
        else:
            raise ValueError("Taux d'intérêt invalide")

    # Méthode statique
    @staticmethod
    def calculer_interets(solde: float, taux: float) -> float:
        return solde * taux

    # Méthodes spéciales
    def __str__(self) -> str:
        return f"Compte de {self.titulaire}: {self._solde:.2f}€"

    def __repr__(self) -> str:
        return f"CompteBancaire(titulaire='{self.titulaire}', solde={self._solde})"

# 2. Utilisation de la classe
compte1 = CompteBancaire("Alice", 1000.0)
print(compte1)  # Utilise __str__

compte1.deposer(500)
print(f"Nouveau solde: {compte1.solde:.2f}€")  # Utilise la propriété

try:
    compte1.retirer(2000)  # Doit lever une exception
except ValueError as e:
    print(f"⚠️ Erreur: {e}")

# 3. Héritage et polymorphisme
class CompteEpargne(CompteBancaire):
    def __init__(self, titulaire: str, solde: float = 0.0, plafond: float = 10000.0):
        super().__init__(titulaire, solde)
        self.plafond = plafond

    # Surcharge de méthode (polymorphisme)
    def retirer(self, montant: float) -> None:
        if 0 < montant <= self._solde and (self._solde - montant) >= 0:
            if self._solde - montant >= 100:  # Solde minimum
                super().retirer(montant)
            else:
                raise ValueError("Solde minimum de 100€ requis")
        else:
            raise ValueError("Montant invalide ou solde insuffisant")

    # Nouvelle méthode spécifique
    def appliquer_interets(self) -> None:
        interets = self.calculer_interets(self._solde, self.taux_interet)
        self.deposer(interets)

# Utilisation de la classe dérivée
compte_epargne = CompteEpargne("Bob", 5000.0)
compte_epargne.appliquer_interets()
print(f"Solde après intérêts: {compte_epargne.solde:.2f}€")

# 4. Encapsulation et propriétés
class Personne:
    def __init__(self, nom: str, age: int):
        self._nom = nom
        self._age = age

    @property
    def nom(self) -> str:
        return self._nom

    @nom.setter
    def nom(self, valeur: str) -> None:
        if not valeur.strip():
            raise ValueError("Le nom ne peut pas être vide")
        self._nom = valeur

    @property
    def age(self) -> int:
        return self._age

    @age.setter
    def age(self, valeur: int) -> None:
        if not 0 <= valeur <= 120:
            raise ValueError("Âge invalide")
        self._age = valeur

# Utilisation des propriétés
p = Personne("Charlie", 30)
print(f"{p.nom} a {p.age} ans")

try:
    p.age = 150  # Doit lever une exception
except ValueError as e:
    print(f"⚠️ {e}")""",
    explanation="""Concepts POO en Python :
1. Classes et instances
2. Attributs (d'instance et de classe)
3. Méthodes (d'instance, de classe, statiques)
4. Encapsulation avec propriétés
5. Héritage et polymorphisme
6. Méthodes spéciales (__str__, __repr__)""",
    notes="""NOTES PRÉSENTATEUR :

    CONCEPTS FONDAMENTAUX DE LA POO :
    - Classe = modèle/plan (blueprint)
    - Objet/Instance = réalisation concrète
    - Attributs = données
    - Méthodes = comportements

    EN PYTHON SPÉCIFIQUEMENT :
    - Tout est objet (même les fonctions et modules)
    - Pas de modificateurs de visibilité stricts (public/private)
    - Conventions: _pour protected, __pour private (name mangling)
    - Méthodes spéciales (__dunder__ methods) pour surcharge d'opérateurs

    TYPES DE MÉTHODES :
    - Méthodes d'instance: prennent self
    - Méthodes de classe: prennent cls, décorateur @classmethod
    - Méthodes statiques: pas de self/cls, décorateur @staticmethod

    PROPRIÉTÉS (@property) :
    - Permettent un contrôle fin sur l'accès aux attributs
    - Getter/setter avec une syntaxe naturelle
    - Utile pour la validation, le calcul à la volée, etc.

    HÉRITAGE :
    - Simple ou multiple
    - super() pour appeler les méthodes parent
    - Méthode resolved order (MRO) pour l'héritage multiple

    POLYMORPHISME :
    - Même interface, comportements différents
    - "Duck typing": si ça marche comme un canard, c'est un canard
    - Pas besoin d'héritage pour le polymorphisme

    CAS D'USAGE EN DATA SCIENCE :
    - Pandas DataFrame est une classe
    - NumPy arrays sont des objets
    - Les visualisations Matplotlib sont des objets
    - Les modèles scikit-learn sont des classes

    BONNES PRATIQUES :
    - Une classe = une responsabilité (principe SRP)
    - Composition > héritage (favoriser la composition d'objets)
    - Noms de classes en PascalCase
    - Noms de méthodes/méthodes en snake_case
    - Documenter avec des docstrings"""
)


# Suite du script - Diapositives 31 à 50

# Diapositive 31 : Approfondissement des listes - Techniques avancées
add_code_slide(
    prs,
    "Techniques avancées avec les listes",
    """# Copie de listes (attention aux références !)
original = [1, 2, 3, 4]
copie1 = original          # Même référence !
copie2 = original.copy()   # Nouvelle liste
copie3 = original[:]       # Autre méthode de copie
original[0] = 99
print("Original:", original)  # [99, 2, 3, 4]
print("Copie1:", copie1)      # [99, 2, 3, 4] - Modifiée !
print("Copie2:", copie2)      # [1, 2, 3, 4] - Intacte
# Listes en compréhension avec conditions
nombres = [1, 2, 3, 4, 5, 6, 7, 8, 9, 10]
pairs = [x for x in nombres if x % 2 == 0]
carrés = [x**2 for x in nombres if x % 2 != 0]
print("Nombres pairs:", pairs)    # [2, 4, 6, 8, 10]
print("Carrés impairs:", carrés)  # [1, 9, 25, 49, 81]
# Fonction zip() pour itérer sur plusieurs listes
noms = ["Alice", "Bob", "Charlie"]
ages = [25, 30, 35]
for nom, age in zip(noms, ages):
    print(f"{nom} a {age} ans")""",
    explanation="Copie profonde vs superficielle et compréhensions avancées",
    notes="""NOTES PRÉSENTATEUR :

    PIÈGE DES RÉFÉRENCES :
    - copie1 = original crée une nouvelle référence vers la même liste
    - Modifier original modifie aussi copie1
    - Utiliser copy() ou [:] pour une vraie copie

    COMPRÉHENSIONS AVANCÉES :
    - Peut inclure des conditions complexes
    - Très utile pour filtrer et transformer des données
    - Syntaxe concise mais puissante

    FONCTION ZIP() :
    - Permet d'itérer sur plusieurs listes en parallèle
    - Très utile pour combiner des données
    - Équivalent à un JOIN en SQL"""
)

# Diapositive 32 : Manipulation de chaînes avancée
add_code_slide(
    prs,
    "Manipulation avancée de chaînes de caractères",
    """# Méthodes utiles des strings
texte = "  Python est un langage Puissant  "
print(texte.strip())            # "Python est un langage Puissant"
print(texte.lower())            # "  python est un langage puissant  "
print(texte.upper())            # "  PYTHON EST UN LANGAGE PUISSANT  "
print(texte.title())            # "  Python Est Un Langage Puissant  "
print(texte.replace("un", "le")) # "  Python est le langage Puissant  "
# Division et jointure
mots = texte.strip().split()
print(mots)                     # ['Python', 'est', 'un', 'langage', 'Puissant']
nouveau_texte = "-".join(mots)
print(nouveau_texte)            # "Python-est-un-langage-Puissant"
# Formatage avancé avec f-strings
nom = "Marie"
age = 28
taille = 1.68
info = f"""
Fiche personnelle:
    Nom: {nom:>10s}
    Âge: {age:03d} ans
    Taille: {taille:.2f}m
"""
print(info)""",
    explanation="Méthodes intégrées et formatage avancé pour le traitement de texte",
    notes="""NOTES PRÉSENTATEUR :

    MÉTHODES UTILES :
    - strip() : supprime les espaces en début/fin
    - lower()/upper() : changement de casse
    - replace() : remplacement de sous-chaînes
    - split()/join() : division et assemblage

    FORMATAGE AVANCÉ :
    - :>10s : aligne à droite sur 10 caractères
    - :03d : nombre entier sur 3 chiffres avec zéros
    - :.2f : nombre flottant avec 2 décimales

    CAS D'USAGE :
    - Nettoyage de données textuelles
    - Génération de rapports formatés
    - Préparation de données pour l'affichage"""
)

# Diapositive 33 : Les ensembles (sets) - Structure puissante
add_code_slide(
    prs,
    "Les ensembles (sets) : Collections uniques non ordonnées",
    """# Création d'ensembles
fruits = {"pomme", "banane", "orange", "pomme"}
print(fruits)  # {"pomme", "banane", "orange"} - doublons supprimés
# Conversion depuis une liste
nombres = [1, 2, 2, 3, 3, 3, 4]
unique = set(nombres)
print(unique)  # {1, 2, 3, 4}
# Opérations sur les ensembles
a = {1, 2, 3, 4}
b = {3, 4, 5, 6}
print("Union:", a | b)          # {1, 2, 3, 4, 5, 6}
print("Intersection:", a & b)    # {3, 4}
print("Différence:", a - b)     # {1, 2}
print("Différence symétrique:", a ^ b)  # {1, 2, 5, 6}
# Ajout et suppression
a.add(5)
a.remove(1)
print("Après modifications:", a)  # {2, 3, 4, 5}
# Test d'appartenance (très rapide)
print("3 dans a ?", 3 in a)      # True
print("1 dans a ?", 1 in a)      # False""",
    explanation="Ensembles = collections non ordonnées d'éléments uniques",
    notes="""NOTES PRÉSENTATEUR :

    PROPRIÉTÉS DES ENSEMBLES :
    - Éléments uniques (pas de doublons)
    - Non ordonnés (pas d'index)
    - Mutables (mais éléments doivent être immuables)
    - Opérations mathématiques ensemblistes

    CAS D'USAGE :
    - Suppression de doublons
    - Tests d'appartenance rapides
    - Opérations ensemblistes (union, intersection)
    - Filtrage de données uniques

    PERFORMANCE :
    - Test d'appartenance en O(1)
    - Très efficace pour grandes collections
    - Base pour certains algorithmes"""
)

# Diapositive 34 : Gestion des fichiers - Lecture/Ecriture
add_code_slide(
    prs,
    "Gestion des fichiers : Lecture et écriture",
    """# Écriture dans un fichier
with open("exemple.txt", "w", encoding="utf-8") as fichier:
    fichier.write("Première ligne\\n")
    fichier.write("Deuxième ligne avec des accents: éèçà\\n")
    fichier.writelines(["Troisième ligne\\n", "Quatrième ligne\\n"])
# Lecture d'un fichier
with open("exemple.txt", "r", encoding="utf-8") as fichier:
    contenu = fichier.read()
    print("Contenu complet:\\n", contenu)
# Lecture ligne par ligne
with open("exemple.txt", "r", encoding="utf-8") as fichier:
    print("\\nLecture ligne par ligne:")
    for num, ligne in enumerate(fichier, 1):
        print(f"Ligne {num}: {ligne.strip()}")
# Gestion des chemins (module pathlib)
from pathlib import Path
chemin = Path("exemple.txt")
print(f"\\nLe fichier existe: {chemin.exists()}")
print(f"Taille: {chemin.stat().st_size} octets")""",
    explanation="Gestion contextuelle des fichiers avec 'with' et encodage UTF-8",
    notes="""NOTES PRÉSENTATEUR :

    BONNES PRATIQUES :
    - Toujours spécifier l'encodage (utf-8)
    - Utiliser 'with' pour gestion automatique
    - pathlib plus moderne que os.path

    MODES D'OUVERTURE :
    - "r" lecture (par défaut)
    - "w" écriture (écrase)
    - "a" ajout (à la fin)
    - "x" création exclusive
    - "+" lecture/écriture

    GESTION DES ERREURS :
    - FileNotFoundError si fichier inexistant
    - PermissionError si droits insuffisants
    - UnicodeDecodeError si encodage incorrect

    CAS D'USAGE :
    - Logs d'application
    - Configuration
    - Import/export de données"""
)

# Diapositive 35 : Gestion des exceptions - Approfondissement
add_code_slide(
    prs,
    "Gestion avancée des exceptions",
    """# Hiérarchie des exceptions
try:
    resultat = 10 / 0
except ZeroDivisionError as e:
    print(f"Erreur de division par zéro: {e}")
except ArithmeticError as e:
    print(f"Erreur arithmétique: {e}")
except Exception as e:
    print(f"Erreur générale: {e}")
# Création d'exceptions personnalisées
class TemperatureError(Exception):
    """Exception pour températures invalides"""
    def __init__(self, temp, message="Température invalide"):
        self.temp = temp
        self.message = message
        super().__init__(f"{message}: {temp}°C")
# Utilisation
def convertir_celsius_fahrenheit(temp_c):
    if temp_c < -273.15:
        raise TemperatureError(temp_c, "Température sous le zéro absolu")
    return temp_c * 9/5 + 32
try:
    print(convertir_celsius_fahrenheit(-300))
except TemperatureError as e:
    print(f"Erreur: {e}")
# Utilisation de else et finally
try:
    fichier = open("fichier_inexistant.txt")
except FileNotFoundError:
    print("Fichier non trouvé")
else:
    print("Fichier ouvert avec succès")
    fichier.close()
finally:
    print("Ce bloc s'exécute toujours")""",
    explanation="Hiérarchie des exceptions et création d'exceptions personnalisées",
    notes="""NOTES PRÉSENTATEUR :

    HIÉRARCHIE DES EXCEPTIONS :
    - Capturer les exceptions spécifiques en premier
    - Exception générale en dernier recours
    - Voir l'arbre d'héritage dans la documentation

    EXCEPTIONS PERSONNALISÉES :
    - Héritent de Exception ou sous-classes
    - Utile pour domaine spécifique
    - Peut inclure des données supplémentaires

    BLOC FINALLY :
    - Toujours exécuté
    - Idéal pour nettoyage (fermeture de fichiers, connexions)
    - Même si une exception est levée

    BONNES PRATIQUES :
    - Ne pas utiliser except: (trop large)
    - Documenter les exceptions levées
    - Utiliser des messages d'erreur clairs"""
)

# Diapositive 36 : Programmation fonctionnelle - Lambda et fonctions d'ordre supérieur
add_code_slide(
    prs,
    "Programmation fonctionnelle : Lambda et fonctions d'ordre supérieur",
    """# Fonctions lambda (anonymes)
carré = lambda x: x ** 2
print(carré(5))  # 25
somme = lambda a, b: a + b
print(somme(3, 4))  # 7
# Utilisation avec sorted()
mots = ["pomme", "banane", "orange", "kiwi", "ananas"]
print(sorted(mots))                     # Tri alphabétique
print(sorted(mots, key=lambda x: len(x))) # Tri par longueur
# Fonctions d'ordre supérieur
def appliquer_fonction(fonction, valeur):
    return fonction(valeur)
print(appliquer_fonction(lambda x: x*3, 5))  # 15
# Map, filter, reduce
nombres = [1, 2, 3, 4, 5]
carrés = list(map(lambda x: x**2, nombres))
pairs = list(filter(lambda x: x % 2 == 0, nombres))
from functools import reduce
produit = reduce(lambda x, y: x * y, nombres)
print("Carrés:", carrés)      # [1, 4, 9, 16, 25]
print("Pairs:", pairs)        # [2, 4]
print("Produit:", produit)    # 120 (1*2*3*4*5)""",
    explanation="Paradigme fonctionnel : fonctions comme objets de première classe",
    notes="""NOTES PRÉSENTATEUR :

    FONCTIONS LAMBDA :
    - Fonctions anonymes en une ligne
    - Syntaxe: lambda [args]: expression
    - Utile pour opérations simples
    - À éviter pour logique complexe

    FONCTIONS D'ORDRE SUPÉRIEUR :
    - Prennent des fonctions en argument
    - Retournent des fonctions
    - map(), filter(), reduce()

    CAS D'USAGE :
    - Tri personnalisé (key=)
    - Transformation de données (map)
    - Filtrage de données (filter)
    - Agrégation (reduce)

    PERFORMANCE :
    - map/filter souvent plus rapides que compréhensions
    - Mais moins lisibles pour débutants
    - reduce moins utilisé (souvent remplacé par boucles)"""
)

# Diapositive 37 : Modules et packages - Organisation du code
add_code_slide(
    prs,
    "Modules et packages : Organisation et réutilisation du code",
    """# Structure d'un package
"""
mon_package/
├── __init__.py        # Initialisation du package
├── module1.py         # Module avec fonctions
├── module2.py         # Autre module
└── sous_package/      # Sous-package
    ├── __init__.py
    └── module3.py
"""
# Contenu de __init__.py
__all__ = ["module1", "module2"]  # Liste des modules exportés
# Importation relative
# Dans module3.py:
from ..module1 import ma_fonction
# Importation absolue
import mon_package.module1
# Utilisation de __name__
if __name__ == "__main__":
    print("Ce code s'exécute seulement si le fichier est lancé directement")
# Installation de packages
# pip install -e .  # Installation en mode éditable""",
    explanation="Architecture modulaire pour projets complexes",
    notes="""NOTES PRÉSENTATEUR :

    STRUCTURE DES PACKAGES :
    - __init__.py marque un dossier comme package
    - Peut être vide ou contenir du code d'initialisation
    - __all__ contrôle ce qui est importé avec from package import *

    IMPORTS :
    - Absolus : depuis la racine du projet
    - Relatifs : avec . (même niveau), .. (niveau supérieur)
    - À privilégier pour éviter les conflits

    BONNES PRATIQUES :
    - Un module = une responsabilité
    - Noms de modules en minuscules
    - Éviter les imports circulaires
    - Documentation avec docstrings

    DISTRIBUTION :
    - setup.py pour la configuration
    - pip install -e pour développement
    - PyPI pour la publication"""
)

# Diapositive 38 : Itérateurs et générateurs
add_code_slide(
    prs,
    "Itérateurs et générateurs : Traitement paresseux des données",
    """# Création d'un itérateur
class Compteur:
    def __init__(self, max):
        self.max = max
        self.current = 0
    def __iter__(self):
        return self
    def __next__(self):
        if self.current >= self.max:
            raise StopIteration
        self.current += 1
        return self.current
compteur = Compteur(3)
for num in compteur:
    print(num)  # 1, 2, 3
# Générateur avec yield
def générateur_nombres(pas=1):
    num = 0
    while True:
        yield num
        num += pas
gen = générateur_nombres(2)
print(next(gen))  # 0
print(next(gen))  # 2
print(next(gen))  # 4
# Expression génératrice
carrés = (x**2 for x in range(5))
for carré in carrés:
    print(carré)  # 0, 1, 4, 9, 16
# Avantages des générateurs
def lire_gros_fichier(nom_fichier):
    with open(nom_fichier) as f:
        for ligne in f:
            yield ligne.strip()""",
    explanation="Traitement mémoire-efficace avec yield et expressions génératrices",
    notes="""NOTES PRÉSENTATEUR :

    ITÉRATEURS VS GÉNÉRATEURS :
    - Itérateur : classe avec __iter__ et __next__
    - Générateur : fonction avec yield
    - Les deux implémentent le protocole d'itération

    AVANTAGES :
    - Traitement paresseux (lazy evaluation)
    - Économie de mémoire
    - Idéal pour grands jeux de données
    - Permet des séquences infinies

    CAS D'USAGE :
    - Traitement de gros fichiers
    - Streams de données
    - Séquences mathématiques infinies
    - Pipelines de traitement

    EXPRESSIONS GÉNÉRATRICES :
    - Syntaxe similaire aux compréhensions de liste
    - Mais utilise () au lieu de []
    - Produit des valeurs à la demande"""
)

# Diapositive 39 : Décorateurs - Métaprogrammation
add_code_slide(
    prs,
    "Décorateurs : Modification dynamique de fonctions",
    """# Décorateur simple
def mon_decorateur(fonction):
    def wrapper():
        print("Avant l'appel de la fonction")
        resultat = fonction()
        print("Après l'appel de la fonction")
        return resultat
    return wrapper
@mon_decorateur
def dire_bonjour():
    print("Bonjour !")
dire_bonjour()
# Décorateur avec arguments
def repetitif(n):
    def decorateur(fonction):
        def wrapper(*args, **kwargs):
            for _ in range(n):
                resultat = fonction(*args, **kwargs)
            return resultat
        return wrapper
    return decorateur
@repetitif(3)
def saluer(nom):
    print(f"Salut {nom}!")
saluer("Alice")  # Affiche 3 fois
# Décorateurs intégrés
@staticmethod
def methode_statique():
    print("Méthode statique")
@classmethod
def methode_classe(cls):
    print(f"Méthode de classe de {cls.__name__}")
# Cache/mémoization
from functools import lru_cache
@lru_cache(maxsize=128)
def fibonacci(n):
    if n < 2:
        return n
    return fibonacci(n-1) + fibonacci(n-2)""",
    explanation="Fonctions qui modifient d'autres fonctions - Puissant mais à utiliser avec parcimonie",
    notes="""NOTES PRÉSENTATEUR :

    CONCEPT DE DÉCORATEUR :
    - Fonction qui prend une fonction en argument
    - Retourne une nouvelle fonction modifiée
    - Syntaxe @decorateur

    CAS D'USAGE :
    - Logging
    - Mesure de performance
    - Vérification d'arguments
    - Cache/mémoization
    - Contrôle d'accès

    DÉCORATEURS INTÉGRÉS :
    - @staticmethod : méthode sans self
    - @classmethod : méthode de classe
    - @property : getter/setter
    - @lru_cache : mémoization

    BONNES PRATIQUES :
    - Garder les décorateurs simples
    - Documenter leur comportement
    - Éviter les effets de bord
    - Utiliser functools.wraps pour préserver les métadonnées"""
)

# Diapositive 40 : Contexte managers - Gestion des ressources
add_code_slide(
    prs,
    "Contexte managers : Gestion automatique des ressources",
    """# Utilisation basique avec 'with'
with open("fichier.txt", "w") as f:
    f.write("Hello World")
# Création d'un contexte manager
class MonContexte:
    def __enter__(self):
        print("Entrée dans le contexte")
        return self
    def __exit__(self, exc_type, exc_val, exc_tb):
        print("Sortie du contexte")
        if exc_type is not None:
            print(f"Exception capturée: {exc_type}")
        return False  # Propager l'exception
with MonContexte() as ctx:
    print("Dans le bloc with")
# Contexte manager avec contextlib
from contextlib import contextmanager
@contextmanager
def mon_contexte():
    print("Setup")
    yield "valeur"
    print("Teardown")
with mon_contexte() as valeur:
    print(f"Valeur reçue: {valeur}")
# Exemple pratique : minuterie
from time import time
@contextmanager
def minuter(nom):
    start = time()
    yield
    elapsed = time() - start
    print(f"{nom} a pris {elapsed:.2f} secondes")""",
    explanation="Protocole pour gestion sûre des ressources avec with",
    notes="""NOTES PRÉSENTATEUR :

    PROTOCOLE CONTEXTE MANAGER :
    - __enter__ : setup, retourne la ressource
    - __exit__ : teardown, gère les exceptions
    - Utilisé avec 'with'

    AVANTAGES :
    - Gestion automatique des ressources
    - Code plus sûr et plus lisible
    - Gestion centralisée des exceptions

    CAS D'USAGE :
    - Ouverture/fermeture de fichiers
    - Connexions réseau/base de données
    - Verrouillage de threads
    - Mesure de performance

    CONTEXTLIB :
    - contextmanager pour créer des CM avec générateurs
    - closing pour ajouter __exit__ à des objets
    - suppress pour ignorer des exceptions

    BONNES PRATIQUES :
    - Toujours libérer les ressources
    - Gérer proprement les exceptions
    - Documenter le comportement"""
)

# Diapositive 41 : Métaclasses - Programmation avancée
add_code_slide(
    prs,
    "Métaclasses : Contrôle de la création des classes",
    """# Métaclasse basique
class MaMeta(type):
    def __new__(cls, name, bases, dct):
        print(f"Création de la classe {name}")
        dct['version'] = 1.0
        return super().__new__(cls, name, bases, dct)
class MaClasse(metaclass=MaMeta):
    pass
print(MaClasse.version)  # 1.0
# Métaclasse pour singleton
class Singleton(type):
    _instances = {}
    def __call__(cls, *args, **kwargs):
        if cls not in cls._instances:
            cls._instances[cls] = super().__call__(*args, **kwargs)
        return cls._instances[cls]
class SingletonClass(metaclass=Singleton):
    pass
a = SingletonClass()
b = SingletonClass()
print(a is b)  # True - Même instance
# Métaclasse pour enregistrement automatique
class PluginMeta(type):
    def __init__(cls, name, bases, dct):
        super().__init__(name, bases, dct)
        if not hasattr(cls, 'plugins'):
            cls.plugins = []
        else:
            cls.plugins.append(cls)
class PluginBase(metaclass=PluginMeta):
    pass
class Plugin1(PluginBase): pass
class Plugin2(PluginBase): pass
print(PluginBase.plugins)  # [<class '__main__.Plugin1'>, <class '__main__.Plugin2'>]""",
    explanation="Contrôle avancé de la création et du comportement des classes",
    notes="""NOTES PRÉSENTATEUR :

    CONCEPT DE MÉTACLASSE :
    - Classe d'une classe (type est la métaclasse par défaut)
    - Contrôle la création des classes
    - Puissant mais complexe

    CAS D'USAGE :
    - Singleton
    - Enregistrement automatique
    - Validation de classes
    - Génération de code
    - Framework ORM

    MÉTACLASSE VS DÉCORATEURS :
    - Métaclasses : contrôle la création des classes
    - Décorateurs : modifie les fonctions/méthodes
    - Préférer les décorateurs quand possible

    BONNES PRATIQUES :
    - Documenter abondamment
    - Éviter sauf nécessité
    - Préférer la composition à l'héritage
    - Tester rigoureusement"""
)

# Diapositive 42 : Programmation asynchrone - async/await
add_code_slide(
    prs,
    "Programmation asynchrone : async/await",
    """# Fonction asynchrone basique
import asyncio
async def dire_bonjour():
    print("Bonjour")
    await asyncio.sleep(1)
    print("Au revoir")
# Exécution d'une coroutine
async def main():
    await dire_bonjour()
# asyncio.run(main())  # À exécuter dans un environnement async
# Exécution concurrentielle
async def compter(n):
    for i in range(n):
        print(f"{i} (compteur {n})")
        await asyncio.sleep(0.1)
async def main_concurrent():
    await asyncio.gather(compter(5), compter(3))
# asyncio.run(main_concurrent())
# Utilisation pratique : requêtes HTTP
import aiohttp
async def fetch_url(url):
    async with aiohttp.ClientSession() as session:
        async with session.get(url) as response:
            return await response.text()
# Exemple avec timeout
async def operation_limitee():
    try:
        async with asyncio.timeout(1):
            await asyncio.sleep(2)
    except TimeoutError:
        print("Opération trop longue !")""",
    explanation="Programmation concurrentielle avec coroutines et event loop",
    notes="""NOTES PRÉSENTATEUR :

    CONCEPTS CLÉS :
    - Coroutine : fonction asynchrone (avec async)
    - await : point de suspension
    - Event loop : orchestre l'exécution
    - Non-bloquant : libère le thread pendant l'attente

    AVANTAGES :
    - Meilleure utilisation des ressources
    - Code plus lisible que les threads
    - Idéal pour I/O bound operations

    CAS D'USAGE :
    - Requêtes HTTP
    - Traitement de streams
    - Serveurs web
    - Bases de données

    BIBLIOTHÈQUES :
    - asyncio : bibliothèque standard
    - aiohttp : requêtes HTTP
    - asyncpg : PostgreSQL
    - aioredis : Redis

    BONNES PRATIQUES :
    - Éviter le code bloquant
    - Gérer les timeouts
    - Limiter la concurrence
    - Tester avec pytest-asyncio"""
)

# Diapositive 43 : Typage statique - Type hints
add_code_slide(
    prs,
    "Typage statique : Type hints et mypy",
    """# Annotation de types basique
def saluer(nom: str) -> str:
    return f"Bonjour {nom}"
# Types complexes
from typing import List, Dict, Tuple, Optional, Union
def traiter_données(
    données: List[Dict[str, Union[int, float]]],
    facteur: Optional[float] = None
) -> Tuple[float, float]:
    # Implementation ici
    return 0.0, 0.0
# Classes avec typage
class Personne:
    def __init__(self, nom: str, age: int):
        self.nom = nom
        self.age = age
    def se_présenter(self) -> str:
        return f"Je m'appelle {self.nom} et j'ai {self.age} ans"
# Typage des variables
age: int = 25
noms: List[str] = ["Alice", "Bob"]
# Union de types (Python 3.10+)
def convertir(valeur: int | str) -> int:
    return int(valeur)
# Utilisation avec mypy
# $ mypy mon_fichier.py  # Vérification statique des types""",
    explanation="Amélioration de la maintenabilité avec annotations de type",
    notes="""NOTES PRÉSENTATEUR :

    AVANTAGES DU TYPAGE :
    - Meilleure documentation
    - Détection précoce d'erreurs
    - Meilleure complétion IDE
    - Facilite la maintenance

    OUTILS :
    - mypy : vérificateur de types statique
    - pyright : alternative de Microsoft
    - pytype : de Google
    - IDEs : PyCharm, VS Code

    TYPES AVANCÉS :
    - List, Dict, Tuple, Set (depuis typing)
    - Optional pour valeurs nulles
    - Union pour plusieurs types
    - TypeVar pour génériques

    BONNES PRATIQUES :
    - Commencer par les fonctions publiques
    - Utiliser graduellement
    - Ne pas sur-annoter
    - Documenter les types complexes"""
)

# Diapositive 44 : Tests unitaires - pytest
add_code_slide(
    prs,
    "Tests unitaires avec pytest",
    """# Structure d'un test
def addition(a: int, b: int) -> int:
    return a + b
def test_addition():
    assert addition(2, 3) == 5
    assert addition(-1, 1) == 0
    assert addition(0, 0) == 0
# Utilisation de fixtures
import pytest
@pytest.fixture
def sample_data():
    return [1, 2, 3, 4, 5]
def test_somme(sample_data):
    assert sum(sample_data) == 15
# Tests paramétrés
@pytest.mark.parametrize("a,b,expected", [
    (1, 2, 3),
    (0, 0, 0),
    (-1, 1, 0)
])
def test_addition_paramétrée(a, b, expected):
    assert addition(a, b) == expected
# Gestion des exceptions
def test_division_par_zero():
    with pytest.raises(ZeroDivisionError):
        1 / 0
# Mocking
from unittest.mock import MagicMock
def test_appel_api():
    mock = MagicMock()
    mock.get.return_value.status_code = 200
    # Utilisation du mock dans le test
# Exécution des tests
# $ pytest mon_fichier.py -v""",
    explanation="Méthodologie de test pour code robuste et maintenable",
    notes="""NOTES PRÉSENTATEUR :

    PRINCIPES DES TESTS :
    - Isoler le code testé
    - Tests déterministes
    - Vérifier un comportement, pas une implémentation
    - Rapides à exécuter

    STRUCTURE :
    - Arrange : préparation
    - Act : action
    - Assert : vérification

    OUTILS :
    - pytest : framework de test
    - coverage : mesure de couverture
    - mock : simulation d'objets
    - hypothesis : tests basés sur propriétés

    BONNES PRATIQUES :
    - Un test = un comportement
    - Noms descriptifs
    - Tests indépendants
    - Exécuter souvent (CI/CD)
    - Cibler 80-90% de couverture"""
)

# Diapositive 45 : Bonnes pratiques - PEP 8 et au-delà
add_content_slide(
    prs,
    "Bonnes pratiques de développement Python",
    """# Conventions de nommage (PEP 8)
variable_snake_case = 42
CONSTANTE_UPPER_CASE = 3.14
def fonction_snake_case():
    pass
class MaClassePascalCase:
    pass
# Structure du code
# 1. Imports standard
# 2. Imports tiers
# 3. Imports locaux
# 4. Constantes
# 5. Classes
# 6. Fonctions
# Documentation
def fonction_bien_documentée(param1: int, param2: str) -> bool:
    \"\"\"
    Description claire de ce que fait la fonction.

    Args:
        param1: Description du premier paramètre
        param2: Description du second paramètre

    Returns:
        bool: Description de la valeur retournée

    Raises:
        ValueError: Si les paramètres sont invalides
    \"\"\"
    # Implémentation
# Outils de qualité
# flake8 : vérification PEP 8
# pylint : analyse statique
# black : formatage automatique
# isort : tri des imports
# Outils de packaging
# poetry : gestion des dépendances
# twine : publication sur PyPI""",
    explanation="Standards et outils pour un code professionnel et maintenable",
    notes="""NOTES PRÉSENTATEUR :

    PEP 8 :
    - 79 caractères par ligne (docstrings: 72)
    - 4 espaces d'indentation
    - Lignes vides pour séparer les fonctions/classes
    - Espaces autour des opérateurs

    DOCUMENTATION :
    - Docstrings pour modules, classes, fonctions
    - Format Google, NumPy ou reStructuredText
    - Exemples dans les docstrings

    OUTILS RECOMMANDÉS :
    - black : formatage automatique
    - isort : organisation des imports
    - flake8/pylint : linting
    - mypy : typage
    - pytest : tests

    BONNES PRATIQUES :
    - Commits atomiques
    - Messages de commit clairs
    - Revue de code
    - Documentation à jour
    - Gestion des dépendances"""
)

# Diapositive 46 : Optimisation des performances
add_content_slide(
    prs,
    "Optimisation des performances en Python",
    """# Profilage du code
import cProfile
def fonction_lente():
    total = 0
    for i in range(1000000):
        total += i
    return total
# cProfile.run('fonction_lente()')
# Optimisations courantes
# 1. Éviter les boucles inutiles
# 2. Utiliser des compréhensions
# 3. Préférer les générateurs pour les gros datasets
# 4. Utiliser des structures de données adaptées
# 5. Minimiser les appels de fonction dans les boucles
# Utilisation de NumPy pour les calculs vectorisés
import numpy as np
data = np.random.rand(1000000)
result = data * 2 + 1  # Opération vectorisée (beaucoup plus rapide)
# Compilation avec Numba
from numba import jit
@jit(nopython=True)
def fonction_optimisée(x):
    return x * 2 + 1
# Utilisation de Cython
# %load_ext cython
# %%cython
# def fonction_cython(int x):
#     return x * 2 + 1
# Parallelisation
from multiprocessing import Pool
with Pool(4) as p:
    results = p.map(fonction_lente, range(4))""",
    explanation="Techniques pour améliorer les performances du code Python",
    notes="""NOTES PRÉSENTATEUR :

    PROFILAGE :
    - Identifier les goulots d'étranglement
    - cProfile pour analyse détaillée
    - timeit pour micro-benchmarks

    OPTIMISATIONS :
    - Algorithmes avant micro-optimisations
    - Structures de données adaptées
    - Éviter les copies inutiles
    - Cache/mémoization

    OUTILS :
    - NumPy : calculs vectorisés
    - Numba : compilation JIT
    - Cython : compilation en C
    - multiprocessing : parallélisme

    BONNES PRATIQUES :
    - Ne pas optimiser prématurément
    - Mesurer avant/après
    - Documenter les optimisations
    - Considérer les trade-offs (lisibilité vs performance)"""
)

# Diapositive 47 : Intégration avec d'autres langages
add_content_slide(
    prs,
    "Intégration de Python avec d'autres langages",
    """# Appel de code C avec ctypes
from ctypes import CDLL, c_int
libc = CDLL("libc.so.6")
print(libc.printf(b"Hello from C!\\n"))  # Appel de fonction C
# Intégration avec C++ via pybind11
# Nécessite compilation séparée
# from my_module import ma_fonction_cpp
# Utilisation de Cython
# fichier.pyx :
# cdef extern from "ma_biblio.h":
#     int ma_fonction_c(int x)
# def wrapper_python(x):
#     return ma_fonction_c(x)
# Appel de Java avec JPype
import jpype
jvm = jpype.startJVM()
java_list = jpype.java.util.ArrayList()
java_list.add("Python")
java_list.add("Java")
print(java_list.size())  # 2
# Intégration avec R via rpy2
import rpy2.robjects as robjects
r = robjects.r
r('x <- c(1, 2, 3)')
print(r['x'][0])  # 1.0
# Appel de code JavaScript avec PyExecJS
import execjs
ctx = execjs.compile("""
    function add(a, b) {
        return a + b;
    }
""")
print(ctx.call("add", 2, 3))  # 5""",
    explanation="Interopérabilité avec d'autres écosystèmes pour étendre les capacités",
    notes="""NOTES PRÉSENTATEUR :

    INTÉGRATION C/CTYPES :
    - Appel direct de bibliothèques C
    - Pas besoin de code intermédiaire
    - Performances natives

    CYTHON :
    - Écriture de code Python avec annotations de type C
    - Compilation en module C
    - Performances proches du C

    JPYPE/RPY2 :
    - Pont entre Python et JVM
    - Intégration avec l'écosystème R
    - Utile pour réutiliser du code existant

    CAS D'USAGE :
    - Réutilisation de bibliothèques existantes
    - Optimisation de parties critiques
    - Intégration dans des systèmes hétérogènes

    OUTILS :
    - ctypes : standard library
    - CFFI : alternative moderne
    - pybind11 : pour C++
    - SWIG : génération de wrappers"""
)

# Diapositive 48 : Déploiement et packaging
add_content_slide(
    prs,
    "Déploiement et packaging d'applications Python",
    """# Structure d'un projet Python
"""
mon_projet/
├── mon_projet/          # Package principal
│   ├── __init__.py
│   ├── module1.py
│   └── module2.py
├── tests/               # Tests unitaires
│   ├── __init__.py
│   └── test_module1.py
├── docs/                # Documentation
├── setup.py             # Configuration du package
├── requirements.txt     # Dépendances
├── README.md            # Documentation
└── pyproject.toml       # Configuration moderne
"""
# Contenu de setup.py
from setuptools import setup, find_packages
setup(
    name="mon_projet",
    version="0.1.0",
    packages=find_packages(),
    install_requires=[
        'numpy>=1.20.0',
        'pandas>=1.3.0'
    ],
    entry_points={
        'console_scripts': [
            'mon_commande=mon_projet.module1:main'
        ]
    }
)
# Création d'un exécutable avec PyInstaller
# pyinstaller --onefile mon_script.py
# Création d'un package wheel
# python setup.py bdist_wheel
# Déploiement sur PyPI
# twine upload dist/*
# Utilisation de Docker
"""
FROM python:3.9-slim
WORKDIR /app
COPY requirements.txt .
RUN pip install -r requirements.txt
COPY . .
CMD ["python", "mon_script.py"]
"""",
    explanation="De la structure du projet au déploiement en production",
    notes="""NOTES PRÉSENTATEUR :

    STRUCTURE STANDARD :
    - Séparation code/source/tests
    - Documentation intégrée
    - Fichiers de configuration clairs

    SETUP.PY :
    - Métadonnées du package
    - Dépendances
    - Points d'entrée pour CLI

    OUTILS DE PACKAGING :
    - setuptools : standard
    - poetry : moderne, gestion des dépendances
    - pipenv : alternative
    - flit : léger

    DÉPLOIEMENT :
    - PyPI pour les bibliothèques
    - Docker pour les applications
    - Serverless (AWS Lambda, etc.)
    - PaaS (Heroku, etc.)

    BONNES PRATIQUES :
    - Versionnement sémantique
    - Documentation complète
    - Tests automatisés
    - CI/CD pipeline
    - Gestion des dépendances"""
)

# Diapositive 49 : Sécurité en Python
add_content_slide(
    prs,
    "Bonnes pratiques de sécurité en Python",
    """# Gestion sécurisée des entrées utilisateur
import re
def valider_email(email: str) -> bool:
    pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\\.[a-zA-Z]{2,}$'
    return bool(re.match(pattern, email))
# Protection contre les injections
import sqlite3
def requete_securisée(db_path, user_id):
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    # Utiliser des paramètres plutôt que de la concaténation
    cursor.execute("SELECT * FROM users WHERE id = ?", (user_id,))
    # Pas : cursor.execute(f"SELECT * FROM users WHERE id = {user_id}")
# Gestion des mots de passe
from passlib.hash import pbkdf2_sha256
mot_de_passe = "mon_mot_de_passe"
hash = pbkdf2_sha256.hash(mot_de_passe)
# Vérification
pbkdf2_sha256.verify(mot_de_passe, hash)  # True
# Sécurité des fichiers
import os
def écrire_fichier_sécurisé(chemin, contenu):
    # Vérifier que le chemin est dans le répertoire autorisé
    if not os.path.abspath(chemin).startswith("/chemin/autorisé"):
        raise ValueError("Chemin non autorisé")
    with open(chemin, "w") as f:
        f.write(contenu)
# Utilisation de HTTPS
import requests
response = requests.get("https://api.example.com", verify=True)  # Vérifie le certificat SSL
# Gestion des secrets
import os
from dotenv import load_dotenv
load_dotenv()  # Charge les variables d'environnement depuis .env
db_password = os.getenv("DB_PASSWORD")  # Jamais en dur dans le code !""",
    explanation="Protection contre les vulnérabilités courantes dans les applications Python",
    notes="""NOTES PRÉSENTATEUR :

    VALIDATION DES ENTRÉES :
    - Toujours valider les données utilisateur
    - Utiliser des expressions régulières
    - Rejeter plutôt que de corriger

    INJECTIONS :
    - SQL : utiliser des paramètres
    - OS : éviter os.system avec entrée utilisateur
    - HTML/JS : échapper les caractères

    MOTS DE PASSE :
    - Jamais en clair dans le code
    - Utiliser des fonctions de hachage (pas MD5/SHA1)
    - pbkdf2, bcrypt, argon2

    GESTION DES SECRETS :
    - Variables d'environnement
    - Fichiers .env (exclus du versioning)
    - Services de gestion de secrets (Vault, etc.)

    BONNES PRATIQUES :
    - Mises à jour régulières
    - Audit des dépendances (safety, dependabot)
    - Principes de moindre privilège
    - Journalisation des activités sensibles"""
)

# Diapositive 50 : Conclusion et ressources
add_content_slide(
    prs,
    "Conclusion et ressources pour aller plus loin",
    """🎯 RÉCAPITULATIF DES COMPÉTENCES ACQUISES :
• Maîtrise de la syntaxe Python et des structures de données
• Compréhension des concepts avancés (POO, décorateurs, métaclasses)
• Capacité à écrire du code robuste et maintenable
• Connaissance des bonnes pratiques et outils professionnels
• Préparation pour l'analyse de données avec NumPy/Pandas
📚 RESSOURCES POUR APPROFONDIR :
• Livres :
  - "Fluent Python" - Luciano Ramalho
  - "Python Cookbook" - David Beazley
  - "Effective Python" - Brett Slatkin
• Sites :
  - Real Python (realpython.com)
  - Python Official Docs (docs.python.org)
  - PyCon Talks (youtube.com/user/pycon101)
• Communautés :
  - Stack Overflow (stackoverflow.com)
  - Python Discord (pythondiscord.com)
  - Meetups locaux (meetup.com)
💡 PROCHAINES ÉTAPES :
• Pratiquer avec des projets personnels
• Explorer les bibliothèques scientifiques (NumPy, Pandas)
• Contribuer à des projets open source
• Participer à des hackathons ou défis de codage
• Rester à jour avec les évolutions du langage""",
    notes="""NOTES PRÉSENTATEUR :

    RÉCAPITULATIF :
    - Faire un tour rapide des concepts clés vus
    - Souligner les compétences acquises
    - Encourager les participants

    RESSOURCES :
    - Recommander des ressources adaptées au niveau
    - Insister sur l'importance de la pratique
    - Encourager la participation communautaire

    PROCHAINES ÉTAPES :
    - Parler des ateliers de l'après-midi
    - Rappeler les objectifs des jours suivants
    - Proposer des idées de projets concrets

    MOTIVATION :
    - Python est un langage en constante évolution
    - Beaucoup d'opportunités professionnelles
    - Communauté très active et accueillante

    CLÔTURE :
    - Remercier les participants
    - Rappeler les horaires de l'après-midi
    - Inviter aux questions finales"""
)

# Enregistrer la présentation mise à jour
output_filename = "Formation_Python_Analyse_Donnees_Jour1_Complete.pptx"
prs.save(output_filename)
print(f"Présentation complète générée et enregistrée sous le nom : {output_filename}")
